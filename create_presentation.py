#!/usr/bin/env python3
"""
Script para criar apresentação PowerPoint do CorrectMe
Versão 2.0 - Com melhorias de design e conteúdo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Cores do tema CorrectMe
PURPLE = RGBColor(168, 85, 247)    # #a855f7
PINK = RGBColor(236, 72, 153)      # #ec4899
DARK_BG = RGBColor(15, 10, 31)     # #0f0a1f
DARK_700 = RGBColor(37, 28, 69)    # #251c45
DARK_600 = RGBColor(55, 42, 92)    # #372a5c
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
LIGHT_GRAY = RGBColor(200, 200, 210)
GREEN = RGBColor(34, 197, 94)
CYAN = RGBColor(34, 211, 238)
BLUE = RGBColor(59, 130, 246)
ORANGE = RGBColor(249, 115, 22)
RED = RGBColor(239, 68, 68)
YELLOW = RGBColor(234, 179, 8)

# Diretório base
BASE_DIR = "/home/user/Correctmemycleverbot"
IMAGES_DIR = os.path.join(BASE_DIR, "images")

# Contador global de páginas
page_number = 0

def add_page_number(slide, prs, number):
    """Adiciona número da página no rodapé"""
    footer_box = slide.shapes.add_textbox(
        Inches(9.2), prs.slide_height - Inches(0.35),
        Inches(0.6), Inches(0.3)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def add_background(slide, prs):
    """Adiciona background escuro padrão"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return bg

def add_title_slide(prs):
    """Slide 1: Título principal com capa atrativa"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Elemento decorativo - gradiente lateral esquerdo
    deco_left = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height
    )
    deco_left.fill.solid()
    deco_left.fill.fore_color.rgb = PURPLE
    deco_left.line.fill.background()

    # Elemento decorativo - linha diagonal
    deco_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(0), Inches(0.08), prs.slide_height
    )
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = PINK
    deco_line.line.fill.background()
    deco_line.rotation = 15

    # Logo
    logo_path = os.path.join(IMAGES_DIR, "LogobrancofundoTransparentev1.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.8), Inches(0.4), height=Inches(1.3))

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CorrectMe"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtítulo atualizado
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(6), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Aprendizagem de Alemão com IA"
    p.font.size = Pt(26)
    p.font.color.rgb = PURPLE

    # Descrição
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(5.5), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Plataforma completa de aprendizado de alemão com inteligência artificial para brasileiros"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

    # Nome do autor
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(5), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dr. Thiago R. S. Pastro"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # Imagem decorativa - Cards representando funcionalidades
    # Card 1 - Correcao
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(0.8), Inches(2.8), Inches(1.3)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = DARK_700
    card1.line.color.rgb = PINK
    card1.line.width = Pt(2)

    card1_title = slide.shapes.add_textbox(Inches(7), Inches(1.0), Inches(2.4), Inches(0.4))
    tf = card1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Correção de Textos"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    card1_icon = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(2.4), Inches(0.5))
    tf = card1_icon.text_frame
    p = tf.paragraphs[0]
    p.text = "IA corrige erros gramaticais"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Card 2 - Conversacao
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.3), Inches(2.8), Inches(1.3)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = DARK_700
    card2.line.color.rgb = PURPLE
    card2.line.width = Pt(2)

    card2_title = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(2.4), Inches(0.4))
    tf = card2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Prática de Conversação"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    card2_icon = slide.shapes.add_textbox(Inches(7), Inches(2.9), Inches(2.4), Inches(0.5))
    tf = card2_icon.text_frame
    p = tf.paragraphs[0]
    p.text = "Fale com personagens IA"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Card 3 - Jogos
    card3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.8), Inches(2.8), Inches(1.3)
    )
    card3.fill.solid()
    card3.fill.fore_color.rgb = DARK_700
    card3.line.color.rgb = GREEN
    card3.line.width = Pt(2)

    card3_title = slide.shapes.add_textbox(Inches(7), Inches(4.0), Inches(2.4), Inches(0.4))
    tf = card3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Jogos e Flashcards"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    card3_icon = slide.shapes.add_textbox(Inches(7), Inches(4.4), Inches(2.4), Inches(0.5))
    tf = card3_icon.text_frame
    p = tf.paragraphs[0]
    p.text = "Aprenda se divertindo"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def add_index_slide(prs):
    """Slide 2: Indice da apresentacao"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Índice"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Linha decorativa
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(1.0), Inches(3), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()

    # Items do indice
    index_items = [
        ("01", "Vantagens do CorrectMe", PINK),
        ("02", "Sistema de Categorização de Erros", PURPLE),
        ("03", "Recursos da Plataforma", CYAN),
        ("04", "Chatbot e Ferramentas IA", GREEN),
        ("05", "Como Funciona", ORANGE),
        ("06", "Análise de Custo", BLUE),
        ("07", "Comparativo e Público Alvo", PINK),
    ]

    y_start = Inches(1.4)
    for i, (num, title, color) in enumerate(index_items):
        y = y_start + i * Inches(0.55)

        # Numero
        num_box = slide.shapes.add_textbox(Inches(2), y, Inches(0.6), Inches(0.45))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color

        # Linha tracejada
        dash_box = slide.shapes.add_textbox(Inches(2.6), y + Inches(0.1), Inches(0.5), Inches(0.3))
        tf = dash_box.text_frame
        p = tf.paragraphs[0]
        p.text = "—"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY

        # Titulo do item
        item_box = slide.shapes.add_textbox(Inches(3.1), y, Inches(5), Inches(0.45))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE

    add_page_number(slide, prs, page_number)

def add_advantages_slide(prs):
    """Slide 3: Vantagens (antes era 'Numeros que Impressionam')"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Vantagens do CorrectMe"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Por que o CorrectMe é a melhor escolha para aprender alemão"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Cards de vantagens com descricao
    advantages = [
        ("+60", "Tópicos de Gramática",
         "Cobertura completa de todos os tópicos gramaticais do alemão, do básico ao avançado", PINK),
        ("+50", "Temas de Conversa",
         "Prática com diferentes personagens e contextos do cotidiano alemão", PURPLE),
        ("5", "Categorias de Erros",
         "Classificação inteligente que ajuda a focar nos seus pontos fracos", BLUE),
        ("24/7", "Disponível Sempre",
         "Acesso ilimitado à plataforma, estude quando e onde quiser", GREEN),
    ]

    card_width = Inches(2.15)
    card_height = Inches(3.2)
    start_x = Inches(0.45)
    y = Inches(1.4)
    gap = Inches(0.2)

    for i, (number, label, desc, color) in enumerate(advantages):
        x = start_x + i * (card_width + gap)

        # Card background com bordas arredondadas
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_700
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Icone circular no topo
        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.65), y + Inches(0.25), Inches(0.85), Inches(0.85)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = color
        icon_circle.line.fill.background()

        # Numero dentro do icone
        num_box = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.35), Inches(0.85), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(22) if len(number) <= 2 else Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.25), card_width - Inches(0.2), Inches(0.6))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Descricao
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.85), card_width - Inches(0.2), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def add_categories_slide(prs):
    """Slide 4: Sistema de Categorizacao de Erros (versao vendedora)"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo mais atrativo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Entenda Exatamente Onde Você Erra"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitulo vendedor
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Nosso sistema exclusivo de categorização identifica e classifica cada erro por cor, \npermitindo que você foque exatamente nos pontos que precisa melhorar"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    categories = [
        ("Declinação", "Genus, Numerus, Kasus - Os 3 pilares da gramática alemã", PINK,
         "Domine os casos e nunca mais erre der/die/das"),
        ("Conjugação", "Tempos verbais e formas conjugadas", PURPLE,
         "Aprenda a conjugar verbos corretamente em qualquer tempo"),
        ("Preposições", "Uso correto com Akkusativ, Dativ e Genitiv", BLUE,
         "Saiba qual preposição usar e com qual caso"),
        ("Sintaxe", "Estrutura e ordem das palavras nas frases", ORANGE,
         "Construa frases com a ordem correta do alemão"),
        ("Vocabulário", "Escolha de palavras e expressões", GREEN,
         "Use as palavras certas para cada contexto"),
    ]

    y_start = Inches(1.7)
    for i, (name, desc, color, benefit) in enumerate(categories):
        y = y_start + i * Inches(0.75)

        # Card de fundo
        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(8.4), Inches(0.65)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = DARK_700
        card_bg.line.fill.background()

        # Barra lateral colorida
        color_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.12), Inches(0.65)
        )
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = color
        color_bar.line.fill.background()

        # Circulo colorido
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.15), Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Nome da categoria
        name_box = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.05), Inches(1.8), Inches(0.35))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Descricao tecnica
        desc_box = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.35), Inches(2.5), Inches(0.3))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY

        # Beneficio
        benefit_box = slide.shapes.add_textbox(Inches(4.3), y + Inches(0.15), Inches(4.7), Inches(0.4))
        tf = benefit_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = benefit
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY

    add_page_number(slide, prs, page_number)

def add_feature_slide_v2(prs, title, subtitle, description, features, image_name, accent_color, image_width=None):
    """Template melhorado para slides de features. image_width em Inches (default 4.7)"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Tag/Subtitulo pequeno
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(4), Inches(0.35))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = accent_color

    # Titulo grande
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(4.5), Inches(0.9))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Descricao
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(4.3), Inches(0.9))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY

    # Features com design moderno (cards pequenos)
    y_start = Inches(2.4)
    for i, feature in enumerate(features):
        y = y_start + i * Inches(0.7)

        # Mini card
        mini_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(4.3), Inches(0.6)
        )
        mini_card.fill.solid()
        mini_card.fill.fore_color.rgb = DARK_700
        mini_card.line.fill.background()

        # Indicador colorido
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.2), Inches(0.2), Inches(0.2)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = accent_color
        indicator.line.fill.background()

        # Texto da feature
        feat_box = slide.shapes.add_textbox(Inches(1.05), y + Inches(0.12), Inches(3.5), Inches(0.4))
        tf = feat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feature
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    # Imagem com moldura arredondada (simulada)
    img_path = os.path.join(IMAGES_DIR, image_name)
    img_w = Inches(image_width) if image_width else Inches(4.7)
    img_x = Inches(5.0) + (Inches(4.7) - img_w) / 2  # Centraliza se menor
    if os.path.exists(img_path):
        try:
            # Adiciona a imagem
            pic = slide.shapes.add_picture(
                img_path, img_x, Inches(0.5), width=img_w
            )
        except Exception as e:
            print(f"Erro ao adicionar imagem {image_name}: {e}")
            # Placeholder quando imagem nao existe
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(0.5), Inches(4.7), Inches(4.3)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = DARK_600
            placeholder.line.color.rgb = accent_color
            placeholder.line.width = Pt(2)

            ph_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(3.7), Inches(0.5))
            tf = ph_text.text_frame
            p = tf.paragraphs[0]
            p.text = f"[{image_name}]"
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
    else:
        # Placeholder quando imagem nao existe
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(0.5), Inches(4.7), Inches(4.3)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = DARK_600
        placeholder.line.color.rgb = accent_color
        placeholder.line.width = Pt(2)

        ph_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(3.7), Inches(0.5))
        tf = ph_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"[{image_name}]"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def add_how_it_works_slide(prs):
    """Slide: Como Funciona - Atualizado"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Como Funciona"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Comece a aprender em 3 passos simples"
    p.font.size = Pt(16)
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Steps atualizados (removido escolher nivel e exercicios personalizados)
    steps = [
        ("1", "Escreva seu texto",
         "Cole ou digite sua redação em alemão. Nosso editor inteligente aceita qualquer texto.", PINK),
        ("2", "Receba a análise",
         "Nossa IA identifica erros, categoriza por tipo e explica cada correção em português.", PURPLE),
        ("3", "Acompanhe sua evolução",
         "Visualize seu progresso no dashboard, pratique com flashcards e melhore continuamente.", GREEN),
    ]

    card_width = Inches(2.9)
    start_x = Inches(0.55)
    y = Inches(1.6)
    gap = Inches(0.25)

    for i, (num, title, desc, color) in enumerate(steps):
        x = start_x + i * (card_width + gap)

        # Card com borda arredondada
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, Inches(3.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_700
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Numero em circulo
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.0), y + Inches(0.3), Inches(0.9), Inches(0.9)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        num_circle.line.fill.background()

        num_box = slide.shapes.add_textbox(x + Inches(1.0), y + Inches(0.4), Inches(0.9), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Titulo do passo
        step_title = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.35), card_width - Inches(0.3), Inches(0.5))
        tf = step_title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Descricao
        step_desc = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.9), card_width - Inches(0.3), Inches(1.2))
        tf = step_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Setas entre os cards
    for i in range(2):
        x = start_x + (i + 1) * (card_width + gap) - Inches(0.15)
        arrow = slide.shapes.add_textbox(x - Inches(0.15), y + Inches(1.5), Inches(0.3), Inches(0.3))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = ">"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PURPLE
        p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def add_cost_analysis_slide(prs):
    """Slide: Analise de Custo"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Análise de Custo"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Investimento acessível para um aprendizado de qualidade"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Sistema de creditos
    credit_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(0.4))
    tf = credit_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Sistema de Créditos"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN

    credit_items = [
        ("100 créditos", "Grátis para novos usuários"),
        ("Correção de texto", "~20 créditos por correção"),
        ("Paráfrase de texto", "~5 créditos por uso"),
        ("Chatbot conversação", "2,5 créditos por frase"),
        ("Chatbot gramática", "5 créditos por tópico"),
        ("Flashcards e jogos", "Gratuito"),
    ]

    y_start = Inches(1.9)
    for i, (item, desc) in enumerate(credit_items):
        y = y_start + i * Inches(0.45)

        # Bullet
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1), Inches(0.15), Inches(0.15)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = CYAN
        bullet.line.fill.background()

        # Item
        item_box = slide.shapes.add_textbox(Inches(1.0), y, Inches(2), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Descricao
        desc_box = slide.shapes.add_textbox(Inches(3.0), y, Inches(2), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # Comparacao com alternativas
    comp_title = slide.shapes.add_textbox(Inches(5.3), Inches(1.4), Inches(4.2), Inches(0.4))
    tf = comp_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Comparado com Alternativas"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    comparisons = [
        ("Professor particular", "R$ 80-150/hora", RED),
        ("Curso presencial", "R$ 300-800/mês", RED),
        ("Apps premium", "R$ 50-100/mês", YELLOW),
        ("CorrectMe", "A partir de R$ 0", GREEN),
    ]

    y_start = Inches(1.9)
    for i, (item, price, color) in enumerate(comparisons):
        y = y_start + i * Inches(0.65)

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), y, Inches(4.2), Inches(0.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_700
        card.line.fill.background()

        # Item
        item_box = slide.shapes.add_textbox(Inches(5.5), y + Inches(0.1), Inches(2.3), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE

        # Preco
        price_box = slide.shapes.add_textbox(Inches(7.8), y + Inches(0.1), Inches(1.5), Inches(0.35))
        tf = price_box.text_frame
        p = tf.paragraphs[0]
        p.text = price
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.RIGHT

    # Destaque final
    highlight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.55), Inches(9), Inches(0.7)
    )
    highlight_box.fill.solid()
    highlight_box.fill.fore_color.rgb = DARK_700
    highlight_box.line.color.rgb = GREEN
    highlight_box.line.width = Pt(2)

    highlight_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(0.4))
    tf = highlight_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Comece grátis com 100 créditos - Sem cartão de crédito necessário!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def add_comparison_slide(prs):
    """Slide: Comparacao com concorrentes"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.55))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Por que escolher o CorrectMe?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Tabela de comparacao
    headers = ["Recurso", "CorrectMe", "LanguageTool", "Scribbr"]
    rows = [
        ["Especializado em alemão", "100%", "Parcial", "Parcial"],
        ["Feedback pedagógico", "Sim", "Não", "Limitado"],
        ["Categorização de erros", "Sim", "Não", "Não"],
        ["Flashcards integrados", "Sim", "Não", "Não"],
        ["Listas de vocabulário", "Sim", "Não", "Não"],
        ["Chatbot para prática", "Sim", "Não", "Não"],
        ["Prática de conversação", "Sim", "Não", "Não"],
    ]

    start_y = Inches(0.8)
    row_height = Inches(0.5)
    col_widths = [Inches(2.8), Inches(2), Inches(2.2), Inches(2)]
    start_x = Inches(0.5)

    # Header background
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, start_x, start_y, Inches(9), row_height
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = DARK_700
    header_bg.line.fill.background()

    # Headers
    x = start_x
    for i, header in enumerate(headers):
        box = slide.shapes.add_textbox(x, start_y + Inches(0.1), col_widths[i], row_height)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PURPLE if i == 1 else WHITE
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        if i == 0:
            tf.margin_left = Inches(0.15)
        x += col_widths[i]

    # Rows
    for row_idx, row in enumerate(rows):
        y = start_y + (row_idx + 1) * row_height

        # Alternating row background
        if row_idx % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, start_x, y, Inches(9), row_height
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(25, 18, 50)
            row_bg.line.fill.background()

        x = start_x
        for col_idx, cell in enumerate(row):
            box = slide.shapes.add_textbox(x, y + Inches(0.1), col_widths[col_idx], row_height)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(11)

            if col_idx == 1:  # CorrectMe column
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif cell == "Não":
                p.font.color.rgb = RED
            elif cell == "Parcial" or cell == "Limitado":
                p.font.color.rgb = YELLOW
            else:
                p.font.color.rgb = GRAY

            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            if col_idx == 0:
                tf.margin_left = Inches(0.15)
            x += col_widths[col_idx]

    add_page_number(slide, prs, page_number)

def add_target_audience_slide(prs):
    """Slide: Publico Alvo"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ideal Para"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    audiences = [
        ("Estudantes", "Preparando para provas Goethe, TestDaF, DSH", PINK),
        ("Profissionais", "Que precisam escrever em alemão no trabalho", PURPLE),
        ("Professores", "Que querem otimizar correções de alunos", BLUE),
        ("Aprendizes", "De qualquer nível, do iniciante ao avançado", GREEN),
    ]

    card_width = Inches(4.2)
    card_height = Inches(1.4)

    positions = [
        (Inches(0.5), Inches(1.3)),
        (Inches(5.3), Inches(1.3)),
        (Inches(0.5), Inches(2.95)),
        (Inches(5.3), Inches(2.95)),
    ]

    for i, ((title, desc, color), (x, y)) in enumerate(zip(audiences, positions)):
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_700
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Icone circular
        icon = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.3), Inches(0.7), Inches(0.7)
        )
        icon.fill.solid()
        icon.fill.fore_color.rgb = color
        icon.line.fill.background()

        # Titulo
        title_box = slide.shapes.add_textbox(x + Inches(1.1), y + Inches(0.25), card_width - Inches(1.3), Inches(0.45))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # Descricao
        desc_box = slide.shapes.add_textbox(x + Inches(1.1), y + Inches(0.7), card_width - Inches(1.3), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY

    add_page_number(slide, prs, page_number)

def add_cta_slide(prs):
    """Slide Final: Call to Action"""
    global page_number
    page_number += 1

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, prs)

    # Elemento decorativo
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = PURPLE
    deco.line.fill.background()

    # Logo
    logo_path = os.path.join(IMAGES_DIR, "LogobrancofundoTransparentev1.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(3.8), Inches(0.5), height=Inches(1.4))

    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Comece a dominar o alemão hoje!"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Experimente grátis - 100 créditos para novos usuários"
    p.font.size = Pt(20)
    p.font.color.rgb = PURPLE
    p.alignment = PP_ALIGN.CENTER

    # URL em destaque
    url_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.7), Inches(5), Inches(0.8)
    )
    url_card.fill.solid()
    url_card.fill.fore_color.rgb = DARK_700
    url_card.line.color.rgb = PINK
    url_card.line.width = Pt(2)

    url_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.85), Inches(5), Inches(0.5))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "correctme.com.br"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PINK
    p.alignment = PP_ALIGN.CENTER

    # Autor
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.4))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dr. Thiago R. S. Pastro"
    p.font.size = Pt(14)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    add_page_number(slide, prs, page_number)

def main():
    global page_number
    page_number = 0

    # Criar apresentacao
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # Slide 1: Titulo
    add_title_slide(prs)

    # Slide 2: Indice
    add_index_slide(prs)

    # Slide 3: Vantagens (antes era "Numeros que Impressionam")
    add_advantages_slide(prs)

    # Slide 4: Categorias de Erros (versao vendedora)
    add_categories_slide(prs)

    # Slide 5: Feature - Correção Inteligente (atualizado)
    add_feature_slide_v2(
        prs,
        title="Feedback detalhado para cada erro",
        subtitle="Correção Inteligente",
        description="Nossa IA identifica erros de declinação, conjugação, sintaxe, preposições e vocabulário, destacando cada um com cores diferentes para facilitar o aprendizado.",
        features=[
            "Correção gramatical e ortográfica em tempo real",
            "Feedback pedagógico que explica o 'porquê' de cada erro",
            "IA treinada para erros comuns de brasileiros",
            "Exportação do texto corrigido"
        ],
        image_name="Screenshot_Composition_english.png",
        accent_color=PINK
    )

    # Slide 6: Feature - Dashboard/Progresso
    add_feature_slide_v2(
        prs,
        title="Visualize sua evolução em tempo real",
        subtitle="Acompanhe seu Progresso",
        description="Dashboard completo com gráficos e estatísticas detalhadas para acompanhar seu desenvolvimento no alemão.",
        features=[
            "Gráfico de distribuição de erros por categoria",
            "Histórico de erros por redação",
            "Estatísticas gerais de performance",
            "Tracking de evolução ao longo do tempo"
        ],
        image_name="Foto resultados.png",
        accent_color=PURPLE
    )

    # Slide 7: Feature - Prática de Conversação (NOVO - com screenshot anexo)
    add_feature_slide_v2(
        prs,
        title="Pratique alemão falado com IA",
        subtitle="Prática de Conversação",
        description="Converse em tempo real com personagens de IA em diferentes contextos profissionais e do dia a dia.",
        features=[
            "10 personagens com diferentes profissões",
            "Reconhecimento de voz em alemão",
            "Análise de erros após a conversa",
            "Prática de escuta e pronúncia"
        ],
        image_name="screenshot_practice_conversation.png",
        accent_color=CYAN
    )

    # Slide 8: Feature - Vocabulário
    add_feature_slide_v2(
        prs,
        title="Organize e pratique seu vocabulário",
        subtitle="Listas de Vocabulário",
        description="Crie listas personalizadas, importe palavras via CSV e classifique por nível de dificuldade.",
        features=[
            "Múltiplas listas organizadas por tema",
            "Importação de CSV",
            "Classificação por dificuldade (vermelho/amarelo/verde)",
            "Tradução automática para português"
        ],
        image_name="Words.png",
        accent_color=GREEN
    )

    # Slide 9: Feature - Flashcards/Jogos (com screenshot anexo)
    add_feature_slide_v2(
        prs,
        title="Pratique de forma interativa e divertida",
        subtitle="Jogos e Flashcards",
        description="Teste seu conhecimento com flashcards interativos e jogos como o jogo da forca alemão.",
        features=[
            "Flashcards automáticos a partir das redações",
            "Jogo da forca com dicas inteligentes",
            "Filtro por nível de dificuldade",
            "Prática de artigos (der/die/das)"
        ],
        image_name="screenshot_flashcards_game.png",
        accent_color=BLUE
    )

    # Slide 10: Feature - Paráfrase (com screenshot anexo)
    add_feature_slide_v2(
        prs,
        title="Reescreva textos em diferentes estilos",
        subtitle="Parafraseador Inteligente",
        description="Transforme seu texto alemão em diferentes estilos: formal, casual, com emojis, simplificado e mais.",
        features=[
            "6 estilos diferentes de reescrita",
            "Mantém o significado original",
            "Ideal para aprender sinônimos",
            "Melhore seu vocabulário ativo"
        ],
        image_name="Parafrasear.png",
        accent_color=ORANGE
    )

    # Slide 11: Feature - Chatbot Correção de Conversação
    add_feature_slide_v2(
        prs,
        title="Pratique escrita com correção em tempo real",
        subtitle="Chatbot de Conversação",
        description="Converse por escrito com o chatbot e receba correções instantâneas de cada frase que você escreve.",
        features=[
            "Correção imediata de cada frase",
            "Feedback detalhado sobre erros",
            "Prática de escrita contextualizada",
            "Apenas 2,5 créditos por frase"
        ],
        image_name="Screenshot_chatbot_correcaoconversa.png",
        accent_color=PINK,
        image_width=2.35  # 50% do tamanho normal
    )

    # Slide 12: Feature - Chatbot Explicação de Gramática
    add_feature_slide_v2(
        prs,
        title="Tire suas dúvidas de gramática",
        subtitle="Chatbot de Gramática",
        description="Pergunte qualquer dúvida sobre gramática alemã e receba explicações detalhadas e exemplos práticos.",
        features=[
            "Explicações em português",
            "Exemplos práticos de uso",
            "Mais de 60 tópicos gramaticais",
            "5 créditos por tópico explicado"
        ],
        image_name="Screenshot_chatbot_explicacaogramatica.png",
        accent_color=PURPLE,
        image_width=2.35  # 50% do tamanho normal
    )

    # Slide 13: Como funciona
    add_how_it_works_slide(prs)

    # Slide 14: Analise de Custo
    add_cost_analysis_slide(prs)

    # Slide 15: Comparacao
    add_comparison_slide(prs)

    # Slide 16: Publico Alvo
    add_target_audience_slide(prs)

    # Slide 17: CTA
    add_cta_slide(prs)

    # Salvar
    output_path = os.path.join(BASE_DIR, "CorrectMe_Apresentacao.pptx")
    prs.save(output_path)
    print(f"Apresentacao salva em: {output_path}")
    print(f"Total de slides: {page_number}")
    return output_path

if __name__ == "__main__":
    main()
